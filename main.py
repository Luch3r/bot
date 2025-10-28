import json
import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont

class PresentationGenerator:
    def __init__(self, json_file):
        self.json_file = json_file
        self.prs = Presentation()
        self.slide_titles = []  
        self.load_data()
    
    def load_data(self):
        """Загрузка данных из JSON файла"""
        with open(self.json_file, 'r', encoding='utf-8') as f:
            self.data = json.load(f)
    
    def set_presentation_properties(self):
        """Установка свойств презентации"""
        if 'title' in self.data['presentation']:
            self.prs.core_properties.title = self.data['presentation']['title']
        if 'author' in self.data['presentation']:
            self.prs.core_properties.author = self.data['presentation']['author']
    
    def add_slide_numbers(self):
        """Добавление номеров слайдов на все слайды"""
        for i, slide in enumerate(self.prs.slides):
            left = Inches(8.5)
            top = Inches(7)
            width = Inches(1)
            height = Inches(0.5)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = str(i + 1)
            if text_frame.paragraphs:
                text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
                text_frame.paragraphs[0].font.size = Pt(12)
    
    def create_table_of_contents(self):
        """Создание слайда с оглавлением"""
        if len(self.slide_titles) <= 1:
            return
            
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        if slide.shapes.title:
            slide.shapes.title.text = "Оглавление"
        
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        
        text_frame.clear()
        for i, title in enumerate(self.slide_titles[1:], 2): 
            p = text_frame.add_paragraph()
            p.text = f"{i-1}. {title}"
            p.level = 0
            p.font.size = Pt(18)
    
    def create_slide(self, slide_data):
        """Создание слайда на основе данных"""
        layout_index = slide_data.get('layout', 1)
        slide_layout = self.prs.slide_layouts[layout_index]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_text = slide_data.get('title', '')
        if title_text:
            self.slide_titles.append(title_text)
        
        if slide.shapes.title:
            slide.shapes.title.text = title_text
        
        self.handle_slide_content(slide, slide_data)
        
        return slide
    
    def handle_slide_content(self, slide, slide_data):
        """Универсальная обработка контента слайда"""
        layout = slide_data.get('layout', 1)
        
        if layout == 0 and 'subtitle' in slide_data:
            self.add_subtitle(slide, slide_data['subtitle'])
        
        if 'content' in slide_data:
            self.add_content_to_slide(slide, slide_data['content'], layout)
        
        if layout == 3:  
            if 'left_content' in slide_data:
                self.add_content_to_slide(slide, slide_data['left_content'], layout, 'left')
            if 'right_content' in slide_data:
                self.add_content_to_slide(slide, slide_data['right_content'], layout, 'right')
        
        if 'images' in slide_data:
            for img_data in slide_data['images']:
                self.add_image(slide, img_data)
    
    def add_subtitle(self, slide, subtitle_text):
        """Добавление подзаголовка на титульный слайд"""
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 15:  
                placeholder.text = subtitle_text
                return
        
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle_text
    
    def add_content_to_slide(self, slide, content_data, layout, position='center'):
        """Добавление контента на слайд"""
        if layout == 3:  
            if position == 'left':
                left, top, width, height = Inches(0.5), Inches(1.5), Inches(4.2), Inches(5)
            else: 
                left, top, width, height = Inches(4.8), Inches(1.5), Inches(4.2), Inches(5)
        else:
            left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(5)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        
        for item in content_data:
            if item['type'] == 'text':
                self.add_text_item(text_frame, item)
            elif item['type'] == 'table':
                self.add_table_to_slide(slide, item)
    
    def add_text_item(self, text_frame, item):
        """Добавление текстового элемента с форматированием"""
        p = text_frame.add_paragraph()
        p.text = item['text']
        p.level = item.get('level', 0)
        
        if 'style' in item:
            self.apply_text_styles(p, item['style'])
    
    def apply_text_styles(self, paragraph, styles):
        """Применение стилей к тексту"""
        font = paragraph.font
        
        if 'bold' in styles and styles['bold']:
            font.bold = True
        if 'italic' in styles and styles['italic']:
            font.italic = True
        if 'size' in styles:
            font.size = Pt(styles['size'])
        if 'color' in styles:
            color = styles['color']
            font.color.rgb = RGBColor(color[0], color[1], color[2])
    
    def add_table_to_slide(self, slide, table_data):
        """Добавление таблицы на слайд"""
        rows = len(table_data['data'])
        cols = len(table_data['data'][0]) if rows > 0 else 0
        
        if rows == 0 or cols == 0:
            return
        
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(0.6 * rows) 
        
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        
        for i, row in enumerate(table_data['data']):
            for j, cell_text in enumerate(row):
                cell = table.cell(i, j)
                cell.text = str(cell_text)
                
                if cell.text_frame.paragraphs:
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                if i == 0 and table_data.get('header', False):
                    if cell.text_frame.paragraphs:
                        paragraph = cell.text_frame.paragraphs[0]
                        paragraph.font.bold = True
                        paragraph.font.size = Pt(14)
    
    def add_image(self, slide, img_data):
        """Добавление изображения"""
        try:
            if os.path.exists(img_data['path']):
                left = Inches(img_data.get('left', 1))
                top = Inches(img_data.get('top', 1))
                width = Inches(img_data.get('width', 4))
                height = Inches(img_data.get('height', 3))
                
                slide.shapes.add_picture(
                    img_data['path'], 
                    left, top, width, height
                )
                print(f"Изображение добавлено: {img_data['path']}")
            else:
                print(f"Файл изображения не найден: {img_data['path']}")
        except Exception as e:
            print(f"Ошибка при добавлении изображения {img_data['path']}: {e}")
    
    def generate(self):
        """Генерация презентации"""
        try:
            self.set_presentation_properties()
            
            for slide_data in self.data['presentation']['slides']:
                self.create_slide(slide_data)
            
            if self.data['presentation'].get('table_of_contents', False):
                self.create_table_of_contents()
            
            self.add_slide_numbers()
            
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            filename = f"presentation_{timestamp}.pptx"
            self.prs.save(filename)
            
            print(f"Презентация успешно создана: {filename}")
            print(f"Количество слайдов: {len(self.prs.slides)}")
            return filename
            
        except Exception as e:
            print(f"Ошибка при создании презентации: {e}")
            import traceback
            traceback.print_exc()
            return None

def create_test_images():
    """Создание тестовых изображений для презентации"""
    
    if not os.path.exists('test_images'):
        os.makedirs('test_images')
    
    img1 = Image.new('RGB', (800, 600), color=(240, 240, 240))
    draw = ImageDraw.Draw(img1)
    
    draw.ellipse([200, 100, 300, 200], outline='blue', width=3)
    draw.ellipse([400, 100, 500, 200], outline='blue', width=3)
    draw.ellipse([500, 250, 600, 350], outline='blue', width=3)
    draw.ellipse([300, 250, 400, 350], outline='blue', width=3)
    draw.ellipse([100, 250, 200, 350], outline='blue', width=3)
    draw.ellipse([300, 400, 400, 500], outline='blue', width=3)
    
    draw.line([250, 200, 350, 250], fill='red', width=2)
    draw.line([450, 200, 550, 250], fill='red', width=2)
    draw.line([250, 200, 150, 250], fill='red', width=2)
    draw.line([350, 350, 350, 400], fill='green', width=2)
    draw.line([550, 350, 350, 400], fill='green', width=2)
    draw.line([150, 350, 350, 400], fill='green', width=2)
    
    try:
        font = ImageFont.truetype("arial.ttf", 24)
    except:
        font = ImageFont.load_default()
    
    draw.text((350, 50), "Нейронная сеть", fill='black', font=font)
    draw.text((150, 150), "Входной слой", fill='blue', font=font)
    draw.text((350, 300), "Скрытый слой", fill='blue', font=font)
    draw.text((320, 450), "Выходной слой", fill='blue', font=font)
    
    img1.save('test_images/ai_brain.jpg')
    
    img2 = Image.new('RGB', (800, 600), color=(255, 255, 255))
    draw = ImageDraw.Draw(img2)
    
    draw.rectangle([50, 50, 200, 200], outline='green', width=3)
    draw.rectangle([300, 50, 450, 200], outline='blue', width=3)
    draw.rectangle([550, 50, 700, 200], outline='red', width=3)
    draw.rectangle([50, 300, 200, 450], outline='purple', width=3)
    draw.rectangle([300, 300, 450, 450], outline='orange', width=3)
    draw.rectangle([550, 300, 700, 450], outline='brown', width=3)
    
    draw.text((70, 210), "Медицина", fill='black', font=font)
    draw.text((320, 210), "Финансы", fill='black', font=font)
    draw.text((570, 210), "Транспорт", fill='black', font=font)
    draw.text((70, 460), "Образование", fill='black', font=font)
    draw.text((310, 460), "Робототехника", fill='black', font=font)
    draw.text((560, 460), "Искусство", fill='black', font=font)
    
    draw.ellipse([100, 100, 150, 150], fill='lightgreen')  
    draw.rectangle([325, 100, 375, 150], fill='lightblue')  
    draw.polygon([(575, 100), (600, 150), (550, 150)], fill='lightcoral')  
    draw.rectangle([100, 325, 150, 375], fill='lavender')  
    draw.rectangle([325, 325, 375, 375], fill='peachpuff')  
    draw.ellipse([575, 325, 625, 375], fill='lightyellow')  
    
    img2.save('test_images/ai_applications.jpg')
    
    print("Тестовые изображения созданы в папке 'test_images/'")

def create_ai_presentation():
    """Создание расширенной презентации об ИИ с изображениями"""
    ai_data = {
        "presentation": {
            "title": "Презентация о искусственном интеллекте",
            "author": "AI Research Team",
            "table_of_contents": True,
            "slides": [
                {
                    "layout": 0,
                    "title": "Искусственный интеллект",
                    "subtitle": "Современные достижения и перспективы"
                },
                {
                    "layout": 1,
                    "title": "Что такое ИИ?",
                    "content": [
                        {
                            "type": "text",
                            "text": "Определение искусственного интеллекта",
                            "level": 0,
                            "style": {
                                "bold": True,
                                "size": 18,
                                "color": [0, 51, 102]
                            }
                        },
                        {
                            "type": "text",
                            "text": "Искусственный интеллект (ИИ) - это область компьютерных наук, занимающаяся созданием машин и систем, способных выполнять задачи, требующие человеческого интеллекта.",
                            "level": 1
                        },
                        {
                            "type": "text",
                            "text": "Основные направления ИИ:",
                            "level": 1,
                            "style": {
                                "bold": True,
                                "italic": True
                            }
                        },
                        {
                            "type": "text",
                            "text": "Машинное обучение",
                            "level": 2
                        },
                        {
                            "type": "text",
                            "text": "Обработка естественного языка",
                            "level": 2
                        },
                        {
                            "type": "text",
                            "text": "Компьютерное зрение",
                            "level": 2
                        },
                        {
                            "type": "text",
                            "text": "Робототехника",
                            "level": 2
                        }
                    ]
                },
                {
                    "layout": 5,
                    "title": "Визуализация ИИ",
                    "images": [
                        {
                            "path": "test_images/ai_brain.jpg",
                            "left": 1,
                            "top": 1.5,
                            "width": 8,
                            "height": 4.5
                        }
                    ],
                    "content": [
                        {
                            "type": "text",
                            "text": "Схематическое представление нейронной сети",
                            "level": 0,
                            "style": {
                                "bold": True,
                                "size": 14,
                                "color": [102, 0, 51]
                            }
                        }
                    ]
                },
                {
                    "layout": 1,
                    "title": "Статистика развития ИИ",
                    "content": [
                        {
                            "type": "table",
                            "header": True,
                            "data": [
                                ["Год", "Инвестиции ($ млрд)", "Количество стартапов", "Доля компаний использующих ИИ"],
                                ["2020", "50.1", "2450", "35%"],
                                ["2021", "68.4", "3120", "47%"],
                                ["2022", "89.7", "3850", "58%"],
                                ["2023", "115.2", "4520", "67%"]
                            ]
                        },
                        {
                            "type": "text",
                            "text": "Данные показывают устойчивый рост рынка ИИ",
                            "level": 0,
                            "style": {
                                "italic": True,
                                "color": [0, 102, 0]
                            }
                        }
                    ]
                },
                {
                    "layout": 3,
                    "title": "Преимущества и вызовы",
                    "left_content": [
                        {
                            "type": "text",
                            "text": "✅ Преимущества ИИ",
                            "level": 0,
                            "style": {
                                "bold": True,
                                "color": [0, 102, 0]
                            }
                        },
                        {
                            "type": "text",
                            "text": "Автоматизация рутинных задач",
                            "level": 1
                        },
                        {
                            "type": "text",
                            "text": "Повышение точности и скорости",
                            "level": 1
                        },
                        {
                            "type": "text",
                            "text": "Персонализация услуг",
                            "level": 1
                        },
                        {
                            "type": "text",
                            "text": "Прогнозирование и анализ",
                            "level": 1
                        },
                        {
                            "type": "text",
                            "text": "Круглосуточная работа",
                            "level": 1
                        }
                    ],
                    "right_content": [
                        {
                            "type": "text",
                            "text": "⚠️ Вызовы и риски",
                            "level": 0,
                            "style": {
                                "bold": True,
                                "color": [153, 0, 0]
                            }
                        },
                        {
                            "type": "text",
                            "text": "Этические вопросы",
                            "level": 1
                        },
                        {
                            "type": "text",
                            "text": "Потеря рабочих мест",
                            "level": 1
                        },
                        {
                            "type": "text",
                            "text": "Безопасность данных",
                            "level": 1
                        },
                        {
                            "type": "text",
                            "text": "Смещение алгоритмов",
                            "level": 1
                        },
                        {
                            "type": "text",
                            "text": "Регуляторные проблемы",
                            "level": 1
                        }
                    ]
                },
                {
                    "layout": 6,
                    "title": "Примеры применения ИИ",
                    "images": [
                        {
                            "path": "test_images/ai_applications.jpg",
                            "left": 0.5,
                            "top": 1,
                            "width": 9,
                            "height": 4
                        }
                    ],
                    "content": [
                        {
                            "type": "text",
                            "text": "Медицина - диагностика заболеваний",
                            "level": 0,
                            "style": {
                                "bold": True,
                                "size": 16,
                                "color": [0, 0, 128]
                            }
                        },
                        {
                            "type": "text",
                            "text": "Финансы - обнаружение мошенничества",
                            "level": 0,
                            "style": {
                                "bold": True,
                                "size": 16,
                                "color": [0, 0, 128]
                            }
                        },
                        {
                            "type": "text",
                            "text": "Транспорт - беспилотные автомобили",
                            "level": 0,
                            "style": {
                                "bold": True,
                                "size": 16,
                                "color": [0, 0, 128]
                            }
                        }
                    ]
                },
                {
                    "layout": 1,
                    "title": "Технологии машинного обучения",
                    "content": [
                        {
                            "type": "text",
                            "text": "Основные типы машинного обучения:",
                            "level": 0,
                            "style": {
                                "bold": True,
                                "size": 16
                            }
                        },
                        {
                            "type": "text",
                            "text": "Обучение с учителем (Supervised Learning)",
                            "level": 1,
                            "style": {
                                "bold": True,
                                "color": [51, 102, 0]
                            }
                        },
                        {
                            "type": "text",
                            "text": "Классификация, регрессия, прогнозирование",
                            "level": 2
                        },
                        {
                            "type": "text",
                            "text": "Обучение без учителя (Unsupervised Learning)",
                            "level": 1,
                            "style": {
                                "bold": True,
                                "color": [102, 51, 0]
                            }
                        },
                        {
                            "type": "text",
                            "text": "Кластеризация, снижение размерности",
                            "level": 2
                        },
                        {
                            "type": "text",
                            "text": "Обучение с подкреплением (Reinforcement Learning)",
                            "level": 1,
                            "style": {
                                "bold": True,
                                "color": [102, 0, 51]
                            }
                        },
                        {
                            "type": "text",
                            "text": "Игровые AI, робототехника, управление системами",
                            "level": 2
                        }
                    ]
                },
                {
                    "layout": 5,
                    "title": "Сравнение алгоритмов ИИ",
                    "content": [
                        {
                            "type": "table",
                            "header": True,
                            "data": [
                                ["Алгоритм", "Точность", "Скорость обучения", "Интерпретируемость", "Область применения"],
                                ["Дерево решений", "85%", "Быстрая", "Высокая", "Классификация"],
                                ["Случайный лес", "92%", "Средняя", "Средняя", "Классификация, Регрессия"],
                                ["Нейронная сеть", "96%", "Медленная", "Низкая", "Компьютерное зрение, NLP"],
                                ["Метод опорных векторов", "88%", "Средняя", "Средняя", "Классификация"],
                                ["K-ближайших соседей", "82%", "Быстрая", "Высокая", "Классификация, Регрессия"]
                            ]
                        }
                    ]
                },
                {
                    "layout": 1,
                    "title": "Будущее искусственного интеллекта",
                    "content": [
                        {
                            "type": "text",
                            "text": "Ключевые тренды на 2024-2030 годы:",
                            "level": 0,
                            "style": {
                                "bold": True,
                                "size": 18,
                                "color": [128, 0, 128]
                            }
                        },
                        {
                            "type": "text",
                            "text": "🤖 Генеративный ИИ и творческие системы",
                            "level": 1,
                            "style": {
                                "bold": True
                            }
                        },
                        {
                            "type": "text",
                            "text": "Создание контента, дизайн, музыка",
                            "level": 2
                        },
                        {
                            "type": "text",
                            "text": "🔬 Научные открытия с помощью ИИ",
                            "level": 1,
                            "style": {
                                "bold": True
                            }
                        },
                        {
                            "type": "text",
                            "text": "Фармацевтика, материаловедение, астрономия",
                            "level": 2
                        },
                        {
                            "type": "text",
                            "text": "🌍 ИИ для устойчивого развития",
                            "level": 1,
                            "style": {
                                "bold": True
                            }
                        },
                        {
                            "type": "text",
                            "text": "Климатическое моделирование, оптимизация энергопотребления",
                            "level": 2
                        },
                        {
                            "type": "text",
                            "text": "⚡ Эффективные и экологичные алгоритмы",
                            "level": 1,
                            "style": {
                                "bold": True
                            }
                        },
                        {
                            "type": "text",
                            "text": "Снижение энергопотребления, edge computing",
                            "level": 2
                        }
                    ]
                },
                {
                    "layout": 0,
                    "title": "Спасибо за внимание!",
                    "subtitle": "Вопросы и обсуждение"
                }
            ]
        }
    }
    
    with open('ai_presentation.json', 'w', encoding='utf-8') as f:
        json.dump(ai_data, f, ensure_ascii=False, indent=2)
    
    return 'ai_presentation.json'

if __name__ == "__main__":
    print("🚀 Запуск генератора презентаций...")
    
    print("📷 Создание тестовых изображений...")
    create_test_images()
    
    print("📄 Создание JSON файла с расширенной презентацией...")
    json_file = create_ai_presentation()
    print(f"✅ Создан файл: {json_file}")
    
    print("🎨 Генерация презентации...")
    generator = PresentationGenerator(json_file)
    result = generator.generate()
    
    if result:
        print("✅ Презентация успешно создана!")
        print(f"📁 Файл: {result}")
    else:
        print("❌ Не удалось создать презентацию")